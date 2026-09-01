"""Extractors for the modern Office formats."""

from __future__ import annotations

from pathlib import Path

from .base import Block, Document, Part, EXTRACTOR_VERSION, source_info

_repairing = {}
EMBEDDED_MEDIA = {
    ".mov",
    ".mp4",
    ".m4v",
    ".avi",
    ".wmv",
    ".mkv",
    ".webm",
    ".m4a",
    ".mp3",
    ".wav",
    ".aac",
    ".wma",
}


def embedded_media_blocks(path: Path) -> list:
    import zipfile

    blocks = []
    try:
        with zipfile.ZipFile(path) as z:
            for info in z.infolist():
                suffix = Path(info.filename).suffix.lower()
                if "/media/" not in info.filename or suffix not in EMBEDDED_MEDIA:
                    continue
                video = suffix not in (".m4a", ".mp3", ".wav", ".aac", ".wma")
                blocks.append(
                    Block(
                        kind="embedded_media",
                        loc={
                            "member": info.filename,
                            "name": Path(info.filename).name,
                            "size_mb": round(info.file_size / 1024 / 1024, 1),
                            "needs_asr": True,
                            "needs_ocr": video,
                            "images": 1 if video else 0,
                        },
                        parts=[],
                    )
                )
    except Exception:
        pass
    return blocks


MAX_SHEET_ROWS = 300

GROUP_MAX_DEPTH = 8


def _walk_shapes(shapes, out: dict, depth: int = 0) -> None:
    """Shapes recursively, groups resolved, sorted by position.

    The old parser walked slide.shapes flat and lost everything inside groups,
    every table and every chart.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        items = sorted(
            shapes,
            key=lambda s: ((getattr(s, "top", None) or 0), (getattr(s, "left", None) or 0)),
        )
    except Exception:
        items = list(shapes)
    for sh in items:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP and depth < GROUP_MAX_DEPTH:
                _walk_shapes(sh.shapes, out, depth + 1)
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                out["images"] += 1
        except Exception:
            pass
        try:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    ph = None
                    try:
                        ph = sh.placeholder_format.type
                    except Exception:
                        pass
                    if ph is not None and "TITLE" in str(ph) and not out["title"]:
                        out["title"] = t
                    else:
                        out["body"].append(t)
        except Exception:
            pass
        try:
            if getattr(sh, "has_table", False):
                out["table"].append(_table_md(sh.table))
        except Exception:
            pass
        try:
            if getattr(sh, "has_chart", False):
                c = _chart_text(sh.chart)
                if c:
                    out["chart"].append(c)
        except Exception:
            pass


def _table_md(table) -> str:

    rows = []
    for r in table.rows:
        rows.append([c.text.strip().replace("\n", " ") for c in r.cells])
    if not rows:
        return ""
    out = ["| " + " | ".join(rows[0]) + " |", "|" + "|".join(["---"] * len(rows[0])) + "|"]
    for r in rows[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def _chart_text(chart) -> str:
    """Data series and categories. Every naive parser throws these away, although
    on consulting slides they often carry the actual message."""
    parts_ = []
    try:
        if chart.has_title and chart.chart_title.text_frame.text.strip():
            parts_.append(chart.chart_title.text_frame.text.strip())
    except Exception:
        pass
    for plot in chart.plots:
        try:
            kats = [str(c) for c in plot.categories if c is not None]
        except Exception:
            kats = []
        for s in plot.series:
            try:
                vals = [v for v in s.values if v is not None]
                paare = ", ".join(f"{k}={v}" for k, v in zip(kats, vals))
                parts_.append(
                    f"{s.name or 'Reihe'}: {paare}" if paare else str(s.name or "")
                )
            except Exception:
                pass
    return "\n".join(t for t in parts_ if t.strip())


def extract_pptx(path: Path) -> Document:

    from pptx import Presentation

    prs = Presentation(str(path))
    doc = Document(
        source=source_info(path), extractor="pptx-native", version=EXTRACTOR_VERSION["pptx"]
    )
    for i, slide in enumerate(prs.slides, 1):
        acc = {"title": "", "body": [], "table": [], "chart": [], "images": 0}
        _walk_shapes(slide.shapes, acc)
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        parts = []
        if acc["title"]:
            parts.append(Part("title", acc["title"]))
        if acc["body"]:
            parts.append(Part("body", "\n".join(acc["body"])))
        for t in acc["table"]:
            if t.strip():
                parts.append(Part("table", t))
        for c in acc["chart"]:
            parts.append(Part("chart", c))
        if notes:
            parts.append(Part("notes", notes))
        doc.blocks.append(
            Block(
                kind="slide",
                loc={"slide": i, "title": acc["title"][:120], "images": acc["images"]},
                parts=parts,
            )
        )
    doc.blocks.extend(embedded_media_blocks(path))
    return doc


def _is_heading(p) -> bool:

    try:
        n = str(p.style.name).lower()
        return n.startswith(("heading", "überschrift", "uberschrift", "title", "title"))
    except Exception:
        return False


def extract_docx(path: Path) -> Document:

    import docx as _docx

    d = _docx.Document(str(path))
    doc = Document(
        source=source_info(path), extractor="docx-native", version=EXTRACTOR_VERSION["docx"]
    )
    current = {"title": "", "text": [], "nr": 1}
    sections = []

    def flush():
        if current["title"] or any(t.strip() for t in current["text"]):
            sections.append(dict(current))
        current["title"] = ""
        current["text"] = []
        current["nr"] += 1

    for p in d.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        if _is_heading(p):
            flush()
            current["title"] = t
        else:
            current["text"].append(t)
    flush()
    for i, a in enumerate(sections, 1):
        parts = []
        if a["title"]:
            parts.append(Part("title", a["title"]))
        if a["text"]:
            parts.append(Part("body", "\n".join(a["text"])))
        if parts:
            doc.blocks.append(
                Block("section", {"section": i, "title": a["title"][:120]}, parts)
            )
    for i, t in enumerate(d.tables, 1):
        md = _table_md(t)
        if md.strip():
            doc.blocks.append(Block("table", {"table": i}, [Part("table", md)]))
    kf = []
    for s in d.sections:
        for teil in (s.header, s.footer):
            try:
                t = "\n".join(p.text.strip() for p in teil.paragraphs if p.text.strip())
                if t and t not in kf:
                    kf.append(t)
            except Exception:
                pass
    if kf:
        doc.blocks.append(
            Block("meta", {"kind": "header-fuss"}, [Part("body", "\n".join(kf))])
        )
    return doc


def _classify(text_zellen: int, zahl_zellen: int) -> str:

    ges = text_zellen + zahl_zellen
    if ges == 0:
        return "empty"
    anteil = text_zellen / ges
    if anteil >= 0.6:
        return "text"
    if anteil >= 0.2:
        return "mixed"
    return "numeric"


def _sheet_block(name: str, rows: list, nr: int) -> Block:
    """rows: lists of raw values. Builds markdown and classifies the sheet."""
    txt = num = 0
    for r in rows:
        for v in r:
            if v is None or str(v).strip() == "":
                continue
            if isinstance(v, (int, float)):
                num += 1
            else:
                txt += 1
    typ = _classify(txt, num)
    rows = [r for r in rows if any(str(v).strip() for v in r if v is not None)]
    rows = rows[:MAX_SHEET_ROWS]
    md = []
    if rows:
        header = [str(v) if v is not None else "" for v in rows[0]]
        md.append("| " + " | ".join(header) + " |")
        md.append("|" + "|".join(["---"] * len(header)) + "|")
        for r in rows[1:]:
            md.append("| " + " | ".join(str(v) if v is not None else "" for v in r) + " |")
    return Block(
        kind="sheet",
        loc={
            "sheet": name,
            "nr": nr,
            "type": typ,
            "rows": len(rows),
            "text_cells": txt,
            "num_cells": num,
        },
        parts=[Part("cell", "\n".join(md))] if md else [],
    )


def extract_xlsx(path: Path) -> Document:

    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    if not wb.worksheets:
        wb.close()
        if not _repairing.get(str(path)):
            _repairing[str(path)] = True
            try:
                from . import legacy

                return legacy.repair_xlsx(path)
            finally:
                _repairing.pop(str(path), None)
    doc = Document(
        source=source_info(path), extractor="xlsx-native", version=EXTRACTOR_VERSION["xlsx"]
    )
    for nr, ws in enumerate(wb.worksheets, 1):
        rows = [
            list(r) for r in ws.iter_rows(max_row=MAX_SHEET_ROWS + 50, values_only=True)
        ]
        doc.blocks.append(_sheet_block(ws.title, rows, nr))
    wb.close()
    return doc


def extract_xls(path: Path) -> Document:

    import xlrd

    try:
        wb = xlrd.open_workbook(str(path))
    except xlrd.XLRDError as e:
        if "encrypt" in str(e).lower():
            d = Document(
                source=source_info(path),
                extractor="xls-xlrd",
                version=EXTRACTOR_VERSION["xls"],
            )
            d.status = "encrypted"
            d.error = "passwortgeschuetzt"
            return d
        raise
    doc = Document(
        source=source_info(path), extractor="xls-xlrd", version=EXTRACTOR_VERSION["xls"]
    )
    for nr, sh in enumerate(wb.sheets(), 1):
        rows = [sh.row_values(i) for i in range(min(sh.nrows, MAX_SHEET_ROWS + 50))]
        doc.blocks.append(_sheet_block(sh.name, rows, nr))
    return doc


def extract_xlsb(path: Path) -> Document:

    from pyxlsb import open_workbook

    doc = Document(
        source=source_info(path), extractor="xlsb-pyxlsb", version=EXTRACTOR_VERSION["xlsb"]
    )
    with open_workbook(str(path)) as wb:
        for nr, name in enumerate(wb.sheets, 1):
            with wb.get_sheet(name) as sh:
                rows = []
                for i, row in enumerate(sh.rows()):
                    if i > MAX_SHEET_ROWS + 50:
                        break
                    rows.append([c.v for c in row])
            doc.blocks.append(_sheet_block(name, rows, nr))
    return doc
